#!/usr/bin/env python3
"""Build a professional 16:9 PowerPoint (.pptx) deck from a deck.json spec.

Usage:
    python build_pptx.py <deck.json> <output.pptx>

The spec describes *content only*; this script owns every design decision
(slide size, palette, fonts, spacing), so decks come out consistent no matter
who writes the spec. See references/pptx-spec.md for the full spec format.

Image paths inside the spec are resolved relative to the deck.json location.
The script exits non-zero and lists every problem it finds (unknown slide
type, missing image, malformed table) so the caller can fix the spec and
re-run. A clean run prints "Wrote <path> (<n> slides)".
"""

import json
import os
import sys

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.util import Emu, Inches, Pt
except ImportError:
    sys.exit("python-pptx is not installed. Run: pip install python-pptx")

# ---------------------------------------------------------------- design ---
SLIDE_W_IN, SLIDE_H_IN = 13.333, 7.5
MARGIN = 0.6                      # left/right margin of content slides
CONTENT_W = SLIDE_W_IN - 2 * MARGIN
CONTENT_TOP = 1.5                 # below the title block
CONTENT_BOTTOM = 6.95             # above the footer

FONT = "Segoe UI"

NAVY = RGBColor(0x1F, 0x38, 0x64)     # primary brand color (matches diagrams)
ACCENT = RGBColor(0x2E, 0x75, 0xB6)   # accent blue
LIGHT = RGBColor(0xED, 0xF2, 0xF8)    # light card / band background
BODY = RGBColor(0x3B, 0x3B, 0x3B)     # body text
MUTED = RGBColor(0x7F, 0x7F, 0x7F)    # captions, footers
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE_BLUE = RGBColor(0xBD, 0xD7, 0xEE)  # subtitle text on navy slides

VALID_TYPES = ("title", "section", "bullets", "two_column", "cards", "stats",
               "image", "image_bullets", "table", "timeline", "closing")


# --------------------------------------------------------------- helpers ---
def textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return box, tf


def next_para(tf, first_used):
    """First call reuses the built-in empty paragraph, then adds new ones."""
    if not first_used[0]:
        first_used[0] = True
        return tf.paragraphs[0]
    return tf.add_paragraph()


def set_text(para, text, size, color, bold=False, italic=False, font=FONT):
    """Write text into a paragraph, honoring **bold** markers inside it."""
    parts = str(text).split("**")
    for i, part in enumerate(parts):
        if part == "":
            continue
        run = para.add_run()
        run.text = part
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold or (i % 2 == 1)
        run.font.italic = italic


def make_bullet(para, level=0, color=None):
    """Give a paragraph a real PowerPoint bullet with hanging indent."""
    p_pr = para._p.get_or_add_pPr()
    indent = Emu(Inches(0.25)).emu
    p_pr.set("marL", str(indent * (level + 1)))
    p_pr.set("indent", str(-indent))
    if color is not None:
        bu_clr = p_pr.makeelement(qn("a:buClr"), {})
        srgb = p_pr.makeelement(qn("a:srgbClr"), {"val": str(color)})
        bu_clr.append(srgb)
        p_pr.append(bu_clr)
    bu_font = p_pr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
    p_pr.append(bu_font)
    char = "•" if level == 0 else "–"     # • then –
    bu_char = p_pr.makeelement(qn("a:buChar"), {"char": char})
    p_pr.append(bu_char)


def solid_rect(slide, left, top, width, height, color,
               shape=MSO_SHAPE.RECTANGLE):
    rect = slide.shapes.add_shape(shape, Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    return rect


def image_size(path):
    """Return (width, height) in pixels; Pillow ships with python-pptx."""
    from PIL import Image
    with Image.open(path) as img:
        return img.size


def fit_picture(slide, path, left, top, max_w, max_h):
    """Place a picture inside a bounding box, preserving aspect ratio."""
    px_w, px_h = image_size(path)
    scale = min(max_w / px_w, max_h / px_h)
    w, h = px_w * scale, px_h * scale
    slide.shapes.add_picture(path, Inches(left + (max_w - w) / 2),
                             Inches(top + (max_h - h) / 2),
                             Inches(w), Inches(h))
    return h


def normalize_bullets(items):
    """Accept plain strings or {"text", "level"} dicts."""
    out = []
    for item in items or []:
        if isinstance(item, dict):
            out.append((str(item.get("text", "")), int(item.get("level", 0))))
        else:
            out.append((str(item), 0))
    return out


def bullet_font_size(bullets, chars_per_line=85):
    """Shrink body text when a slide carries a lot of it."""
    lines = sum(max(1, (len(t) + chars_per_line - 1) // chars_per_line)
                for t, _ in bullets)
    if lines <= 7:
        return 18
    if lines <= 10:
        return 16
    return 14


def write_bullets(tf, bullets, size, first_used):
    for text, level in bullets:
        para = next_para(tf, first_used)
        para.space_after = Pt(max(6, size - 8))
        para.line_spacing = 1.0
        set_text(para, text, size - (2 if level else 0),
                 BODY if level == 0 else MUTED)
        make_bullet(para, level, color="2E75B6" if level == 0 else "7F7F7F")


# ---------------------------------------------------- shared slide pieces ---
def content_header(slide, data):
    """Title + accent rule used by every light content slide."""
    title = data.get("title", "")
    _, tf = textbox(slide, MARGIN, 0.42, CONTENT_W, 0.75)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    set_text(para, title, 27 if len(title) <= 55 else 23, NAVY, bold=True)
    solid_rect(slide, MARGIN, 1.22, 1.6, 0.045, ACCENT)
    if data.get("subtitle"):
        _, tf = textbox(slide, MARGIN + 1.8, 1.08, CONTENT_W - 1.8, 0.3)
        set_text(tf.paragraphs[0], data["subtitle"], 13, MUTED, italic=True)


def footer(slide, deck, number):
    if deck.get("footer"):
        _, tf = textbox(slide, MARGIN, 7.08, 8.0, 0.3)
        set_text(tf.paragraphs[0], deck["footer"], 9, MUTED)
    _, tf = textbox(slide, SLIDE_W_IN - MARGIN - 0.8, 7.08, 0.8, 0.3)
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    set_text(tf.paragraphs[0], str(number), 9, MUTED)


def dark_background(slide):
    solid_rect(slide, 0, 0, SLIDE_W_IN, SLIDE_H_IN, NAVY)
    solid_rect(slide, 0, SLIDE_H_IN - 0.12, SLIDE_W_IN, 0.12, ACCENT)


# ----------------------------------------------------------- slide types ---
def build_title(slide, data, deck):
    dark_background(slide)
    solid_rect(slide, MARGIN + 0.15, 2.25, 1.1, 0.05, ACCENT)
    _, tf = textbox(slide, MARGIN + 0.15, 2.5, 11.9, 1.7)
    para = tf.paragraphs[0]
    set_text(para, data.get("title", deck.get("title", "")), 40, WHITE,
             bold=True)
    if data.get("subtitle"):
        para = tf.add_paragraph()
        para.space_before = Pt(14)
        set_text(para, data["subtitle"], 20, PALE_BLUE)
    meta = [m for m in (data.get("meta"), data.get("author"),
                        data.get("date")) if m]
    if meta:
        _, tf = textbox(slide, MARGIN + 0.15, 6.3, 11.9, 0.7)
        para = tf.paragraphs[0]
        set_text(para, "   |   ".join(str(m) for m in meta), 13, PALE_BLUE)


def build_section(slide, data, deck):
    dark_background(slide)
    if data.get("number"):
        _, tf = textbox(slide, MARGIN + 0.15, 2.0, 11.9, 1.3)
        set_text(tf.paragraphs[0], f"{data['number']:02d}"
                 if isinstance(data["number"], int) else str(data["number"]),
                 60, ACCENT, bold=True)
    _, tf = textbox(slide, MARGIN + 0.15, 3.35, 11.9, 1.6)
    set_text(tf.paragraphs[0], data.get("title", ""), 34, WHITE, bold=True)
    if data.get("subtitle"):
        para = tf.add_paragraph()
        para.space_before = Pt(10)
        set_text(para, data["subtitle"], 17, PALE_BLUE)


def build_bullets(slide, data, deck):
    content_header(slide, data)
    bullets = normalize_bullets(data.get("bullets"))
    _, tf = textbox(slide, MARGIN, CONTENT_TOP + 0.15, CONTENT_W,
                    CONTENT_BOTTOM - CONTENT_TOP - 0.15)
    write_bullets(tf, bullets, bullet_font_size(bullets), [False])


def build_two_column(slide, data, deck):
    content_header(slide, data)
    columns = data.get("columns") or []
    col_w = (CONTENT_W - 0.5) / 2
    for i, col in enumerate(columns[:2]):
        left = MARGIN + i * (col_w + 0.5)
        top = CONTENT_TOP + 0.15
        if col.get("heading"):
            solid_rect(slide, left, top, col_w, 0.5, LIGHT)
            _, tf = textbox(slide, left + 0.15, top, col_w - 0.3, 0.5)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            set_text(tf.paragraphs[0], col["heading"], 16, NAVY, bold=True)
            top += 0.7
        bullets = normalize_bullets(col.get("bullets"))
        _, tf = textbox(slide, left, top, col_w, CONTENT_BOTTOM - top)
        write_bullets(tf, bullets,
                      min(16, bullet_font_size(bullets, 45)), [False])


def build_cards(slide, data, deck):
    content_header(slide, data)
    cards = (data.get("cards") or [])[:6]
    per_row = 2 if len(cards) <= 4 else 3
    if len(cards) == 3:
        per_row = 3
    rows = (len(cards) + per_row - 1) // per_row
    gap = 0.3
    card_w = (CONTENT_W - (per_row - 1) * gap) / per_row
    card_h = min(2.4, (CONTENT_BOTTOM - CONTENT_TOP - (rows - 1) * gap) / rows)
    for i, card in enumerate(cards):
        row, col = divmod(i, per_row)
        left = MARGIN + col * (card_w + gap)
        top = CONTENT_TOP + 0.1 + row * (card_h + gap)
        rect = solid_rect(slide, left, top, card_w, card_h, LIGHT,
                          MSO_SHAPE.ROUNDED_RECTANGLE)
        rect.adjustments[0] = 0.055
        solid_rect(slide, left, top + 0.12, 0.05, card_h - 0.24, ACCENT)
        tf = rect.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = tf.margin_right = Inches(0.22)
        tf.margin_top = Inches(0.16)
        para = tf.paragraphs[0]
        set_text(para, card.get("title", ""), 15, NAVY, bold=True)
        para.space_after = Pt(6)
        if card.get("text"):
            para = tf.add_paragraph()
            para.line_spacing = 1.05
            set_text(para, card["text"],
                     12 if len(card["text"]) <= 170 else 11, BODY)


def build_stats(slide, data, deck):
    content_header(slide, data)
    stats = (data.get("stats") or [])[:5]
    gap = 0.3
    card_w = (CONTENT_W - (len(stats) - 1) * gap) / max(len(stats), 1)
    for i, stat in enumerate(stats):
        left = MARGIN + i * (card_w + gap)
        rect = solid_rect(slide, left, 2.2, card_w, 2.1, LIGHT,
                          MSO_SHAPE.ROUNDED_RECTANGLE)
        rect.adjustments[0] = 0.07
        tf = rect.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.1)
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        value = str(stat.get("value", ""))
        set_text(para, value, 36 if len(value) <= 8 else 26, ACCENT, bold=True)
        para = tf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        para.space_before = Pt(6)
        set_text(para, stat.get("label", ""), 12.5, BODY)
    bullets = normalize_bullets(data.get("bullets"))
    if bullets:
        _, tf = textbox(slide, MARGIN, 4.75, CONTENT_W, CONTENT_BOTTOM - 4.75)
        write_bullets(tf, bullets, 15, [False])


def build_image(slide, data, deck, spec_dir):
    content_header(slide, data)
    path = resolve_image(data, spec_dir)
    caption_h = 0.45 if data.get("caption") else 0
    box_h = CONTENT_BOTTOM - CONTENT_TOP - 0.1 - caption_h
    fit_picture(slide, path, MARGIN, CONTENT_TOP + 0.05, CONTENT_W, box_h)
    if data.get("caption"):
        _, tf = textbox(slide, MARGIN, CONTENT_BOTTOM - 0.4, CONTENT_W, 0.35)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        set_text(tf.paragraphs[0], data["caption"], 11.5, MUTED, italic=True)


def build_image_bullets(slide, data, deck, spec_dir):
    content_header(slide, data)
    path = resolve_image(data, spec_dir)
    text_w = 4.3
    img_w = CONTENT_W - text_w - 0.4
    img_on_left = data.get("image_side", "right") == "left"
    img_x = MARGIN if img_on_left else MARGIN + text_w + 0.4
    txt_x = MARGIN + img_w + 0.4 if img_on_left else MARGIN
    fit_picture(slide, path, img_x, CONTENT_TOP + 0.1, img_w,
                CONTENT_BOTTOM - CONTENT_TOP - 0.2)
    bullets = normalize_bullets(data.get("bullets"))
    _, tf = textbox(slide, txt_x, CONTENT_TOP + 0.2, text_w,
                    CONTENT_BOTTOM - CONTENT_TOP - 0.2)
    write_bullets(tf, bullets, min(15, bullet_font_size(bullets, 40)),
                  [False])


def build_table(slide, data, deck):
    content_header(slide, data)
    columns = data.get("columns") or []
    rows = data.get("rows") or []
    if not columns or not rows:
        raise ValueError("table slide needs non-empty 'columns' and 'rows'")
    for r_i, row in enumerate(rows):
        if len(row) != len(columns):
            raise ValueError(f"table row {r_i + 1} has {len(row)} cells, "
                             f"expected {len(columns)}")
    height = min(0.5 + 0.42 * len(rows), CONTENT_BOTTOM - CONTENT_TOP - 0.1)
    shape = slide.shapes.add_table(len(rows) + 1, len(columns),
                                   Inches(MARGIN), Inches(CONTENT_TOP + 0.1),
                                   Inches(CONTENT_W), Inches(height))
    table = shape.table
    for c_i, name in enumerate(columns):
        cell = table.cell(0, c_i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.text_frame.word_wrap = True
        para = cell.text_frame.paragraphs[0]
        set_text(para, name, 13, WHITE, bold=True)
    for r_i, row in enumerate(rows):
        for c_i, value in enumerate(row):
            cell = table.cell(r_i + 1, c_i)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_i % 2 else LIGHT
            cell.text_frame.word_wrap = True
            para = cell.text_frame.paragraphs[0]
            set_text(para, value, 12, BODY)


def build_timeline(slide, data, deck):
    content_header(slide, data)
    steps = (data.get("steps") or [])[:6]
    if not steps:
        raise ValueError("timeline slide needs a non-empty 'steps' list")
    gap = 0.12
    step_w = (CONTENT_W - (len(steps) - 1) * gap) / len(steps)
    for i, step in enumerate(steps):
        left = MARGIN + i * (step_w + gap)
        shape_kind = MSO_SHAPE.PENTAGON if i == 0 else MSO_SHAPE.CHEVRON
        rect = solid_rect(slide, left, 2.35, step_w, 0.85,
                          NAVY if i % 2 == 0 else ACCENT, shape_kind)
        tf = rect.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.12)
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        set_text(para, step.get("label", ""), 13, WHITE, bold=True)
        if step.get("text"):
            _, tf = textbox(slide, left + 0.05, 3.45, step_w - 0.1, 2.6)
            para = tf.paragraphs[0]
            para.line_spacing = 1.05
            set_text(para, step["text"], 11.5, BODY)


def build_closing(slide, data, deck):
    dark_background(slide)
    _, tf = textbox(slide, MARGIN + 0.15, 2.7, 11.9, 1.4)
    set_text(tf.paragraphs[0], data.get("title", "Thank You"), 36, WHITE,
             bold=True)
    if data.get("subtitle"):
        para = tf.add_paragraph()
        para.space_before = Pt(12)
        set_text(para, data["subtitle"], 17, PALE_BLUE)
    lines = data.get("lines") or []
    if lines:
        _, tf = textbox(slide, MARGIN + 0.15, 5.0, 11.9, 1.5)
        first = [False]
        for line in lines:
            para = next_para(tf, first)
            para.space_after = Pt(5)
            set_text(para, line, 13, PALE_BLUE)


# ------------------------------------------------------------- machinery ---
def resolve_image(data, spec_dir):
    raw = data.get("image", "")
    if not raw:
        raise ValueError("slide is missing its 'image' path")
    path = raw if os.path.isabs(raw) else os.path.join(spec_dir, raw)
    path = os.path.normpath(path)
    if not os.path.isfile(path):
        raise FileNotFoundError(f"image not found: {path}")
    return path


DARK_TYPES = ("title", "section", "closing")


def build_deck(deck, spec_dir, out_path):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]

    errors = []
    slides = deck.get("slides") or []
    if not slides:
        sys.exit("deck.json has no slides")

    for index, data in enumerate(slides, start=1):
        kind = data.get("type", "")
        slide = prs.slides.add_slide(blank)
        try:
            if kind == "title":
                build_title(slide, data, deck)
            elif kind == "section":
                build_section(slide, data, deck)
            elif kind == "bullets":
                build_bullets(slide, data, deck)
            elif kind == "two_column":
                build_two_column(slide, data, deck)
            elif kind == "cards":
                build_cards(slide, data, deck)
            elif kind == "stats":
                build_stats(slide, data, deck)
            elif kind == "image":
                build_image(slide, data, deck, spec_dir)
            elif kind == "image_bullets":
                build_image_bullets(slide, data, deck, spec_dir)
            elif kind == "table":
                build_table(slide, data, deck)
            elif kind == "timeline":
                build_timeline(slide, data, deck)
            elif kind == "closing":
                build_closing(slide, data, deck)
            else:
                raise ValueError(
                    f"unknown slide type '{kind}' - valid types: "
                    + ", ".join(VALID_TYPES))
        except Exception as exc:
            errors.append(f"slide {index} ({kind or 'no type'}, "
                          f"'{data.get('title', '')}'): {exc}")
            continue
        if kind not in DARK_TYPES:
            footer(slide, deck, index)
        if data.get("notes"):
            slide.notes_slide.notes_text_frame.text = str(data["notes"])

    if errors:
        print("FAILED - fix deck.json and re-run:")
        for err in errors:
            print(f"  {err}")
        sys.exit(1)

    prs.save(out_path)
    print(f"Wrote {out_path} ({len(slides)} slides)")


def main():
    if len(sys.argv) != 3:
        sys.exit("Usage: python build_pptx.py <deck.json> <output.pptx>")
    spec_path, out_path = sys.argv[1], sys.argv[2]
    if not os.path.isfile(spec_path):
        sys.exit(f"Spec not found: {spec_path}")
    with open(spec_path, encoding="utf-8") as f:
        try:
            deck = json.load(f)
        except json.JSONDecodeError as exc:
            sys.exit(f"deck.json is not valid JSON: {exc}")
    build_deck(deck, os.path.dirname(os.path.abspath(spec_path)),
               out_path)


if __name__ == "__main__":
    main()
